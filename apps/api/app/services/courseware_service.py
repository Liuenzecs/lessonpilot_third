"""课件大纲生成与 PPTX 导出服务。

从教案结构化内容生成课堂教学课件。
"""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

from app.models import Task
from app.schemas.content import (
    DocumentContent,
    LessonPlanContent,
    TeachingProcessStep,
    is_lesson_plan,
)

# ---------------------------------------------------------------------------
# 配色（与 Word 导出一致的中式现代风）
# ---------------------------------------------------------------------------

_COLOR_INK = "2c2c2c"
_COLOR_STONE = "3a7ca5"
_COLOR_MUTED = "8c8c8c"
_COLOR_LIGHT_BG = "f7f5f0"

_SLIDE_WIDTH = Cm(33.867)  # 16:9
_SLIDE_HEIGHT = Cm(19.05)

# ---------------------------------------------------------------------------
# 课件大纲数据模型
# ---------------------------------------------------------------------------


@dataclass
class Slide:
    """单张课件幻灯片大纲。"""

    index: int
    slide_type: str  # title | objectives | key_points | teaching_step | questions | summary | homework
    title: str
    bullet_points: list[str] = field(default_factory=list)
    speaker_notes: str = ""
    step: TeachingProcessStep | None = None


# ---------------------------------------------------------------------------
# 教案 → 课件大纲映射
# ---------------------------------------------------------------------------


def generate_slide_outline(task: Task, content: LessonPlanContent) -> list[Slide]:
    """从教案内容生成课件大纲。"""
    slides: list[Slide] = []
    idx = 0

    # 1. 标题页
    slides.append(_build_title_slide(task, content, idx))
    idx += 1

    # 2. 教学目标
    if content.objectives_status == "confirmed" and content.objectives:
        slides.append(_build_objectives_slide(content, idx))
        idx += 1

    # 3. 教学重难点
    kp = content.key_points
    if content.key_points_status == "confirmed" and (kp.key_points or kp.difficulties):
        slides.append(_build_key_points_slide(content, idx))
        idx += 1

    # 4. 教学过程 — 每个环节一页
    if content.teaching_process_status == "confirmed" and content.teaching_process:
        for step in content.teaching_process:
            slides.append(_build_teaching_step_slide(step, idx))
            idx += 1

    # 5. 课堂提问汇总（从已确认的教学过程中提取）
    questions = _extract_questions(content) if content.teaching_process_status == "confirmed" else []
    if questions:
        slides.append(_build_questions_slide(questions, idx))
        idx += 1

    # 6. 板书 / 总结
    if content.board_design_status == "confirmed" and content.board_design:
        slides.append(_build_summary_slide(content, idx))
        idx += 1

    # 7. 课后作业
    slides.append(_build_homework_slide(idx))
    idx += 1

    return slides


def _build_title_slide(task: Task, content: LessonPlanContent, idx: int) -> Slide:
    header = content.header
    category_labels = {"new": "新授课", "review": "复习课", "exercise": "练习课", "comprehensive": "综合课"}
    category = category_labels.get(task.lesson_category, task.lesson_category)
    return Slide(
        index=idx,
        slide_type="title",
        title=header.title or task.title,
        bullet_points=[
            f"学科：{task.subject}　年级：{task.grade}　课时：第{task.class_hour}课时",
            f"课型：{category}",
        ],
        speaker_notes="",
    )


def _build_objectives_slide(content: LessonPlanContent, idx: int) -> Slide:
    dim_labels = {"knowledge": "知识与技能", "ability": "过程与方法", "emotion": "情感态度与价值观"}
    return Slide(
        index=idx,
        slide_type="objectives",
        title="教学目标",
        bullet_points=[f"【{dim_labels.get(o.dimension, o.dimension)}】{o.content}" for o in content.objectives],
        speaker_notes="",
    )


def _build_key_points_slide(content: LessonPlanContent, idx: int) -> Slide:
    kp = content.key_points
    points: list[str] = []
    if kp.key_points:
        points.append("教学重点：")
        points.extend(f"• {p}" for p in kp.key_points)
    if kp.difficulties:
        if points:
            points.append("")
        points.append("教学难点：")
        points.extend(f"• {d}" for d in kp.difficulties)
    return Slide(
        index=idx,
        slide_type="key_points",
        title="教学重难点",
        bullet_points=points,
        speaker_notes="",
    )


def _build_teaching_step_slide(step: TeachingProcessStep, idx: int) -> Slide:
    title = f"{step.phase}" if step.phase else f"环节{idx}"
    if step.duration:
        title += f"（{step.duration}分钟）"
    return Slide(
        index=idx,
        slide_type="teaching_step",
        title=title,
        bullet_points=_split_smart(step.teacher_activity),
        speaker_notes=f"设计意图：{step.design_intent}\n\n学生活动：{step.student_activity}",
        step=step,
    )


def _extract_questions(content: LessonPlanContent) -> list[str]:
    """从教学过程中提取课堂提问。"""
    questions: list[str] = []
    for step in content.teaching_process:
        for line in step.teacher_activity.replace("\n", "。").split("。"):
            stripped = line.strip()
            has_q = "?" in stripped or "？" in stripped or stripped.startswith("提问") or stripped.startswith("问：")
            if stripped and has_q:
                if len(stripped) > 3:
                    questions.append(stripped)
        for line in step.student_activity.replace("\n", "。").split("。"):
            stripped = line.strip()
            if stripped and ("讨论" in stripped or "思考" in stripped or "回答" in stripped):
                if len(stripped) > 3 and stripped not in questions:
                    questions.append(stripped)
    return questions[:10]


def _build_questions_slide(questions: list[str], idx: int) -> Slide:
    return Slide(
        index=idx,
        slide_type="questions",
        title="课堂互动",
        bullet_points=[f"Q{i + 1}: {q}" for i, q in enumerate(questions)],
        speaker_notes="",
    )


def _build_summary_slide(content: LessonPlanContent, idx: int) -> Slide:
    return Slide(
        index=idx,
        slide_type="summary",
        title="课堂总结",
        bullet_points=_split_smart(content.board_design),
        speaker_notes="",
    )


def _build_homework_slide(idx: int) -> Slide:
    return Slide(
        index=idx,
        slide_type="homework",
        title="课后作业",
        bullet_points=["（请根据教学实际情况布置课后作业）"],
        speaker_notes="",
    )


def _split_smart(text: str) -> list[str]:
    """智能分点：按换行、句号、分号拆分。"""
    if not text:
        return []
    parts = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    result: list[str] = []
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if len(part) > 80:
            sub = part.replace("。", "。\n").replace("；", "；\n").split("\n")
            for s in sub:
                s = s.strip()
                if s:
                    result.append(s)
        else:
            result.append(part)
    return result


# ---------------------------------------------------------------------------
# PPTX 构建
# ---------------------------------------------------------------------------


def build_pptx(task: Task, content: DocumentContent, template_spec: dict | None = None) -> bytes:
    """从教案内容生成 PPTX 文件字节。"""
    if not is_lesson_plan(content):
        raise ValueError("PPTX export is only supported for lesson plans")

    slides_outline = generate_slide_outline(task, content)
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    # 使用空白版式
    blank_layout = prs.slide_layouts[6]  # blank

    for slide_data in slides_outline:
        _add_slide(prs, blank_layout, slide_data)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_slide(prs: Presentation, layout, slide_data: Slide) -> None:
    slide = prs.slides.add_slide(layout)

    # 标题
    title_box = slide.shapes.add_textbox(Cm(2), Cm(1.5), Cm(29.867), Cm(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.title
    p.alignment = PP_ALIGN.LEFT
    _set_run_font(p, "微软雅黑", Pt(28), bold=True, color=_COLOR_STONE)

    # 正文要点
    if slide_data.bullet_points:
        body_box = slide.shapes.add_textbox(Cm(3), Cm(4.5), Cm(27.867), Cm(12))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        for i, point in enumerate(slide_data.bullet_points):
            if i == 0:
                bp = body_tf.paragraphs[0]
            else:
                bp = body_tf.add_paragraph()
            bp.text = point
            bp.space_after = Pt(8)
            bp.line_spacing = 1.5
            _set_run_font(bp, "宋体", Pt(18), color=_COLOR_INK)

    # 演讲者备注
    if slide_data.speaker_notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data.speaker_notes


def _set_run_font(para, name: str, size, *, bold: bool = False, color: str | None = None) -> None:
    """设置段落中 run 的字体属性。"""
    if not para.runs:
        para.add_run()
    run = para.runs[0]
    run.font.name = name
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor.from_string(color)

    # 设置东亚字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = run._r.makeelement(qn("a:ea"), {})
        rPr.insert(0, ea)
    ea.set("typeface", name)
